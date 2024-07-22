from flask import Flask, request, jsonify
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
import tempfile
import uuid
from dotenv import load_dotenv
import pickle
import firebase_admin
from firebase_admin import credentials, storage, firestore
from flask_cors import CORS

cred = credentials.Certificate("config/god-rage-aaca3-firebase-adminsdk-8o3uh-97a23a3858.json")
firebase_admin.initialize_app(cred, {
    'storageBucket': 'god-rage-aaca3.appspot.com'
})
db = firestore.client()
bucket = storage.bucket()

app = Flask(__name__)
CORS(app)
load_dotenv()
os.getenv("PAID_API_KEY")
genai.configure(api_key=os.getenv("PAID_API_KEY"))

sessions = {}

def get_pdf_text(pdf_path):
    text = ""
    pdf_reader = PdfReader(pdf_path)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_ppt_text(ppt_path):
    text = ""
    presentation = Presentation(ppt_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_doc_text(doc_path):
    text = ""
    doc = Document(doc_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1500,
        chunk_overlap=300,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatGoogleGenerativeAI(model="gemini-1.5-pro-latest", temperature=0.5)
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

@app.route('/create_session', methods=['POST'])
def create_session():
    if 'files' not in request.files:
        return jsonify({"error": "No files provided"}), 400

    files = request.files.getlist('files')
    raw_text = ""
    chat_name = request.form['chat_name']
    for file in files:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            file.save(temp_file.name)
            if file.filename.endswith('.pdf'):
                raw_text += get_pdf_text(temp_file.name)
            elif file.filename.endswith('.ppt') or file.filename.endswith('.pptx'):
                raw_text += get_ppt_text(temp_file.name)
            elif file.filename.endswith('.doc') or file.filename.endswith('.docx'):
                raw_text += get_doc_text(temp_file.name)

    session_id = str(uuid.uuid4())
    text_chunks = get_text_chunks(raw_text)
    vectorstore = get_vectorstore(text_chunks)
    conversation_chain = get_conversation_chain(vectorstore)

    sessions[session_id] = {
        "conversation_chain": conversation_chain,
        "chat_history": [],
        "text_chunks": text_chunks
    }

    chunk_size = 100  # Number of chunks per document
    chunk_urls = []
    for i in range(0, len(text_chunks), chunk_size):
        chunk_part = text_chunks[i:i + chunk_size]
        chunk_blob_name = f'{session_id}/chunk_{i // chunk_size}.pkl'
        chunk_blob = bucket.blob(chunk_blob_name)
        chunk_blob.upload_from_string(pickle.dumps(chunk_part))
        chunk_urls.append(chunk_blob_name)

    session_data = {
        'id': session_id,
        'chunk_urls': chunk_urls,
        'chat_history': [],
        'name': chat_name,
    }
    db.collection('sessions').document(session_id).set(session_data)
    
    print("Session created:", session_id)
    return jsonify(session_data), 200

@app.route('/ask_question', methods=['POST'])
def ask_question():
    data = request.json
    session_id = data.get('session_id')
    user_question = data.get('question')

    if not session_id or not user_question:
        return jsonify({"error": "Invalid session_id or question"}), 400

    session = sessions.get(session_id)
    if not session:
        session_doc = db.collection('sessions').document(session_id).get()
        if not session_doc.exists:
            return jsonify({"error": "Session not found"}), 404

        session_data = session_doc.to_dict()
        chunk_urls = session_data.get('chunk_urls', [])

        text_chunks = []
        for chunk_blob_name in chunk_urls:
            chunk_blob = bucket.blob(chunk_blob_name)
            chunk_data = chunk_blob.download_as_string()
            text_chunks.extend(pickle.loads(chunk_data))

        vectorstore = get_vectorstore(text_chunks)
        conversation_chain = get_conversation_chain(vectorstore)

        session = {
            "conversation_chain": conversation_chain,
            "chat_history": session_data.get('chat_history', []),
            "text_chunks": text_chunks
        }
        sessions[session_id] = session
    else:
        conversation_chain = session["conversation_chain"]

    response = conversation_chain({'question': user_question})
    chat_history = response['chat_history']

    messages = []
    for i, message in enumerate(chat_history):
        messages.append({
            "id": str(uuid.uuid4()),  # Add a unique identifier
            "role": "user" if i % 2 == 0 else "bot",
            "content": message.content
        })

    new_message = {
        "id": str(uuid.uuid4()), 
        "isUser": False,
        "message": response['answer']
    }

    db.collection('sessions').document(session_id).update({'chat_history': firestore.ArrayUnion([new_message])})

    session["chat_history"].extend(messages)
    return jsonify({"chat_history": messages}), 200

@app.route('/get_chat_history', methods=['POST'])
def get_chat_history():
    data = request.json
    session_id = data.get('session_id')

    if not session_id:
        return jsonify({"error": "Invalid session_id"}), 400

    session = sessions.get(session_id)

    if not session:
        return jsonify({"error": "Session not found"}), 404

    return jsonify({"chat_history": session['chat_history']}), 200

if __name__ == '__main__':
    app.run(debug=True)
